#!/usr/bin/env python3
"""Generate a minimal PPTX sample with run-level hyperlink text.

Output (from repo root):
  samples/pptx/pptx_hyperlink_basic.pptx

Dependencies:
  pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = REPO_ROOT / "samples" / "pptx"
OUT_FILE = OUT_DIR / "pptx_hyperlink_basic.pptx"


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    slide.shapes.title.text = "PPTX Hyperlink Basic"

    body = slide.placeholders[1].text_frame
    body.clear()

    p = body.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Project Home: "

    run2 = p.add_run()
    run2.text = "https://example.com/project"
    run2.hyperlink.address = "https://example.com/project"

    prs.save(OUT_FILE)
    print(f"[ok] wrote: {OUT_FILE}")


if __name__ == "__main__":
    main()
