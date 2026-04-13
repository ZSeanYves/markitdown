#!/usr/bin/env python3
"""Generate a minimal PPTX sample: single image + single caption-like text.

Output (from repo root):
  samples/pptx/pptx_image_single_caption_like.pptx

Dependencies:
  pip install python-pptx Pillow
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = REPO_ROOT / "samples" / "pptx"
OUT_FILE = OUT_DIR / "pptx_image_single_caption_like.pptx"


def build_temp_image() -> Path:
    """Create a temporary placeholder PNG and return its path."""
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp_path = Path(tmp.name)
    tmp.close()

    img = Image.new("RGB", (960, 540), color=(245, 248, 252))
    draw = ImageDraw.Draw(img)
    draw.rectangle((40, 40, 920, 500), outline=(80, 110, 160), width=6)
    draw.text((70, 70), "Sample Figure", fill=(40, 60, 90))
    img.save(tmp_path)
    return tmp_path


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    img_path = build_temp_image()
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # One single image.
        slide.shapes.add_picture(str(img_path), Inches(1.2), Inches(1.0), width=Inches(8.0))

        # One single caption-like text directly below image.
        caption_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(8.0), Inches(0.6))
        tf = caption_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Figure 1. Revenue trend"
        p.font.size = Pt(16)

        prs.save(OUT_FILE)
        print(f"[ok] wrote: {OUT_FILE}")
    finally:
        try:
            img_path.unlink(missing_ok=True)
        except Exception:
            pass


if __name__ == "__main__":
    main()
