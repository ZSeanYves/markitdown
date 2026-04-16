#!/usr/bin/env python3
"""
Generate binary regression samples for hyperlink coverage.

This script intentionally keeps generated .docx/.pptx files out of git history.
"""

from __future__ import annotations

import argparse
import os
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET


ROOT = Path(__file__).resolve().parent
DOCX_DIR = ROOT / "docx"
PPTX_DIR = ROOT / "pptx"


def _ensure_imports() -> tuple[object, object]:
    try:
        import docx  # type: ignore
        import pptx  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise SystemExit(
            "Missing dependency. Please install: pip install python-docx python-pptx\n"
            f"Original error: {exc}"
        )
    return docx, pptx


def _write_docx_missing_rel(docx_mod: object, force: bool) -> None:
    out_path = DOCX_DIR / "docx_hyperlink_missing_rel_negative.docx"
    if out_path.exists() and not force:
        return

    Document = docx_mod.Document  # type: ignore[attr-defined]
    doc = Document()
    p = doc.add_paragraph("Visit ")

    from docx.oxml import OxmlElement  # type: ignore
    from docx.oxml.ns import qn  # type: ignore

    part = doc.part
    rel_type = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
    )
    r_id = part.relate_to("https://example.com/missing", rel_type, is_external=True)

    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    r = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "Broken Link"
    r.append(t)
    hyperlink.append(r)
    p._p.append(hyperlink)
    p.add_run(" please.")

    tmp_path = out_path.with_suffix(".tmp.docx")
    doc.save(tmp_path)

    # Remove hyperlink relationship to create missing-rel negative case.
    with zipfile.ZipFile(tmp_path, "r") as zin, zipfile.ZipFile(
        out_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/_rels/document.xml.rels":
                root = ET.fromstring(data)
                for rel in list(root):
                    if rel.attrib.get("Type", "").endswith("/hyperlink"):
                        root.remove(rel)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            zout.writestr(item, data)

    tmp_path.unlink(missing_ok=True)


def _write_pptx_shape_samples(pptx_mod: object, force: bool) -> None:
    Presentation = pptx_mod.Presentation  # type: ignore[attr-defined]

    basic_path = PPTX_DIR / "pptx_hyperlink_shape_basic.pptx"
    if force or not basic_path.exists():
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Shape Level Link"
        title.click_action.hyperlink.address = "https://example.com/shape"
        prs.save(str(basic_path))

    neg_path = PPTX_DIR / "pptx_hyperlink_shape_negative_empty_target.pptx"
    if force or not neg_path.exists():
        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[5])
        title2 = slide2.shapes.title
        title2.text = "Shape Link Empty Target"
        title2.click_action.hyperlink.address = "https://example.com/will-be-empty"
        prs2.save(str(neg_path))

        with zipfile.ZipFile(neg_path, "r") as zin:
            files = {i.filename: zin.read(i.filename) for i in zin.infolist()}

        rels_name = "ppt/slides/_rels/slide1.xml.rels"
        root = ET.fromstring(files[rels_name])
        for rel in root:
            if rel.attrib.get("Type", "").endswith("/hyperlink"):
                rel.set("Target", "")
        files[rels_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

        with zipfile.ZipFile(neg_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for name, data in files.items():
                zout.writestr(name, data)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--force", action="store_true")
    args = parser.parse_args()

    DOCX_DIR.mkdir(parents=True, exist_ok=True)
    PPTX_DIR.mkdir(parents=True, exist_ok=True)

    docx_mod, pptx_mod = _ensure_imports()
    _write_docx_missing_rel(docx_mod, force=args.force)
    _write_pptx_shape_samples(pptx_mod, force=args.force)
    print("Generated hyperlink binary samples.")


if __name__ == "__main__":
    main()
