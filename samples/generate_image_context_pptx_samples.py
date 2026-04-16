#!/usr/bin/env python3
"""
Generate PPTX binaries for image-context regression cases.

Outputs:
  - samples/pptx/pptx_image_caption_basic.pptx
  - samples/pptx/pptx_image_caption_near_basic.pptx
  - samples/pptx/pptx_image_multiple_caption_ambiguous_negative.pptx
"""

from __future__ import annotations

from pathlib import Path


def require_deps():
    try:
        from pptx import Presentation  # noqa: F401
        from pptx.util import Inches  # noqa: F401
    except Exception as exc:  # pragma: no cover - generator runtime guard
        raise SystemExit(
            "python-pptx is required to generate PPTX image-context samples.\n"
            "Install with: python3 -m pip install python-pptx"
        ) from exc


def generate(root: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    img = root / "samples" / "image" / "img" / "img_red.jpg"
    if not img.exists():
        raise SystemExit(f"image not found: {img}")

    out_dir = root / "samples" / "pptx"
    out_dir.mkdir(parents=True, exist_ok=True)

    # single image + caption-like text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
    box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(5.5), Inches(0.8))
    box.text_frame.text = "Figure 1. Basic caption"
    prs.save(out_dir / "pptx_image_caption_basic.pptx")

    # single image + nearby short text (fallback path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(4))
    box = slide.shapes.add_textbox(Inches(1.1), Inches(5.2), Inches(5.5), Inches(0.8))
    box.text_frame.text = "Revenue trend chart"
    prs.save(out_dir / "pptx_image_caption_near_basic.pptx")

    # ambiguous multi-image + multi-caption scene
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img), Inches(0.7), Inches(1), width=Inches(3))
    slide.shapes.add_picture(str(img), Inches(6), Inches(1), width=Inches(3))
    b1 = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(3.5), Inches(0.8))
    b1.text_frame.text = "Figure 1. Left sample"
    b2 = slide.shapes.add_textbox(Inches(6), Inches(4.6), Inches(3.5), Inches(0.8))
    b2.text_frame.text = "Figure 2. Right sample"
    prs.save(out_dir / "pptx_image_multiple_caption_ambiguous_negative.pptx")


def main() -> None:
    require_deps()
    root = Path(__file__).resolve().parent.parent
    generate(root)
    print("generated PPTX image-context samples")


if __name__ == "__main__":
    main()
