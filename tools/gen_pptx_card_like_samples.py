#!/usr/bin/env python3
"""Generate card-like PPTX regression samples locally.

Note: this script is optional local tooling and is not invoked in CI.

Dependencies:
  pip install python-pptx
"""

import argparse
from pathlib import Path
import sys

try:
    from pptx import Presentation
    from pptx.util import Pt
except ModuleNotFoundError:  # pragma: no cover - env-dependent
    Presentation = None
    Pt = None

ROOT = Path(__file__).resolve().parents[1]

EMU = int


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def add_textbox(
    slide,
    left: EMU,
    top: EMU,
    width: EMU,
    height: EMU,
    text: str,
    font_size: int = 20,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    return tb


def add_title(slide, text: str) -> None:
    add_textbox(
        slide,
        left=700000,
        top=180000,
        width=7800000,
        height=500000,
        text=text,
        font_size=32,
    )


def add_card_pair(
    slide,
    left: EMU,
    top: EMU,
    title: str,
    description: str,
    card_width: EMU = 2300000,
    title_height: EMU = 420000,
    description_height: EMU = 420000,
    inner_gap: EMU = 10000,
    title_size: int = 22,
    description_size: int = 18,
) -> None:
    add_textbox(
        slide,
        left=left,
        top=top,
        width=card_width,
        height=title_height,
        text=title,
        font_size=title_size,
    )
    add_textbox(
        slide,
        left=left,
        top=top + title_height + inner_gap,
        width=card_width,
        height=description_height,
        text=description,
        font_size=description_size,
    )


def build_pptx_card_pairs_basic(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Capability Cards")

    card_positions = [
        (900000, 1300000, "Search", "Fast retrieval"),
        (4200000, 1300000, "Vision", "Image parsing"),
        (900000, 2800000, "Speech", "Audio input"),
        (4200000, 2800000, "Agents", "Task execution"),
    ]

    for left, top, title, description in card_positions:
        add_card_pair(slide, left, top, title, description)

    filename = "pptx_card_pairs_basic.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_card_pairs_with_side_note(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Deployment Modes")

    card_positions = [
        (900000, 1300000, "Cloud", "Managed hosting"),
        (4200000, 1300000, "On-prem", "Customer environment"),
        (900000, 2800000, "Hybrid", "Split workloads"),
        (4200000, 2800000, "Edge", "Low-latency inference"),
    ]

    for left, top, title, description in card_positions:
        add_card_pair(slide, left, top, title, description)

    add_textbox(
        slide,
        left=6750000,
        top=3650000,
        width=2200000,
        height=420000,
        text="Internal preview only.",
        font_size=16,
    )

    filename = "pptx_card_pairs_with_side_note.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_negative_dense_keyword_wall(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Model Topics")

    words = [
        "Search",
        "Ranking",
        "Safety",
        "Vision",
        "Speech",
        "Agents",
        "Memory",
        "Planning",
        "Tooling",
    ]

    col_x = [1100000, 3500000, 5900000]
    row_y = [1400000, 2100000, 2800000]

    index = 0
    for y in row_y:
        for x in col_x:
            add_textbox(
                slide,
                left=x,
                top=y,
                width=1650000,
                height=360000,
                text=words[index],
                font_size=20,
            )
            index += 1

    filename = "pptx_negative_dense_keyword_wall.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_card_pairs_two_rows_three_cols(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Feature Cards")

    cards = [
        (700000, 1250000, "Search", "Fast retrieval"),
        (3300000, 1250000, "Vision", "Image parsing"),
        (5900000, 1250000, "Speech", "Audio input"),
        (700000, 2600000, "Agents", "Task execution"),
        (3300000, 2600000, "Memory", "Session state"),
        (5900000, 2600000, "Planning", "Multi-step reasoning"),
    ]

    for left, top, title, description in cards:
        add_card_pair(
            slide,
            left=left,
            top=top,
            title=title,
            description=description,
            card_width=2200000,
            title_height=360000,
            description_height=360000,
            inner_gap=6000,
            title_size=21,
            description_size=17,
        )

    filename = "pptx_card_pairs_two_rows_three_cols.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_card_pairs_two_groups(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Platform Options")

    left_group = [
        (900000, 1450000, "Cloud", "Managed hosting"),
        (900000, 2800000, "On-prem", "Customer environment"),
    ]
    right_group = [
        (5200000, 1450000, "Hybrid", "Split workloads"),
        (5200000, 2800000, "Edge", "Low-latency inference"),
    ]

    for left, top, title, description in left_group + right_group:
        add_card_pair(
            slide,
            left=left,
            top=top,
            title=title,
            description=description,
            card_width=2500000,
            title_height=380000,
            description_height=380000,
            inner_gap=8000,
            title_size=22,
            description_size=18,
        )

    filename = "pptx_card_pairs_two_groups.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_negative_caption_scatter(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Annotations")

    labels = [
        (1450000, 1650000, "Search"),
        (4550000, 1800000, "Vision"),
        (3000000, 2350000, "Latency"),
        (1700000, 3050000, "Recall"),
        (5200000, 3150000, "Edge"),
        (3600000, 3650000, "Draft"),
    ]

    for left, top, text in labels:
        add_textbox(
            slide,
            left=left,
            top=top,
            width=1600000,
            height=340000,
            text=text,
            font_size=20,
        )

    filename = "pptx_negative_caption_scatter.pptx"
    prs.save(output_dir / filename)
    return filename


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate selected PPTX card-like regression samples.",
    )
    parser.add_argument(
        "output_dir",
        help="Directory to write generated .pptx files (example: samples/pptx).",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()
    if Presentation is None or Pt is None:
        print("Missing dependency: python-pptx. Install with: pip install python-pptx")
        sys.exit(1)

    out_dir = Path(args.output_dir)
    if not out_dir.is_absolute():
        out_dir = (ROOT / out_dir).resolve()

    ensure_dir(out_dir)

    generated = [
        build_pptx_card_pairs_basic(out_dir),
        build_pptx_card_pairs_with_side_note(out_dir),
        build_pptx_negative_dense_keyword_wall(out_dir),
        build_pptx_card_pairs_two_rows_three_cols(out_dir),
        build_pptx_card_pairs_two_groups(out_dir),
        build_pptx_negative_caption_scatter(out_dir),
    ]

    for name in generated:
        print(f"Generated: {name}")
