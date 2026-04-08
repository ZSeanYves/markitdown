#!/usr/bin/env python3
"""Generate caption-like / note-cluster PPTX samples locally.

Note: this script is local tooling and is not invoked in CI.

Dependencies:
  pip install python-pptx
"""

import argparse
from pathlib import Path
import sys

try:
    from pptx import Presentation
    from pptx.util import Pt
except ModuleNotFoundError:  # pragma: no cover - env-dependent
    Presentation = None
    Pt = None

ROOT = Path(__file__).resolve().parents[1]
EMU = int


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def add_textbox(
    slide,
    left: EMU,
    top: EMU,
    width: EMU,
    height: EMU,
    text: str,
    font_size: int = 20,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    return tb


def add_title(slide, text: str) -> None:
    add_textbox(
        slide,
        left=700000,
        top=180000,
        width=7800000,
        height=500000,
        text=text,
        font_size=32,
    )


def build_pptx_callout_blocks_basic(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Release Notes")

    callouts = [
        (850000, 1550000, "Security", "Secure by default"),
        (3400000, 1700000, "Sharing", "Team access controls"),
        (5950000, 1500000, "Export", "Markdown and PDF output"),
    ]

    for left, top, heading, description in callouts:
        add_textbox(
            slide,
            left=left,
            top=top,
            width=2250000,
            height=360000,
            text=heading,
            font_size=23,
        )
        add_textbox(
            slide,
            left=left,
            top=top + 390000,
            width=2350000,
            height=360000,
            text=description,
            font_size=18,
        )

    filename = "pptx_callout_blocks_basic.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_callout_blocks_row_jitter(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Callout Row Jitter")

    callouts = [
        (800000, 1450000, 1835000, "Title1", "Desc1"),
        (3350000, 1465000, 1860000, "Title2", "Desc2"),
        (5900000, 1435000, 1840000, "Title3", "Desc3"),
    ]
    desc_y_jitter = [390000, 405000, 375000]

    for i, (left, top, width, title, desc) in enumerate(callouts):
        add_textbox(
            slide,
            left=left,
            top=top,
            width=width,
            height=340000,
            text=title,
            font_size=22,
        )
        add_textbox(
            slide,
            left=left,
            top=top + desc_y_jitter[i],
            width=width + 140000,
            height=360000,
            text=desc,
            font_size=18,
        )

    filename = "pptx_callout_blocks_row_jitter.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_callout_blocks_mixed_widths(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Callout Mixed Widths")

    callouts = [
        (700000, 1500000, 1450000, 1880000, "Narrow", "Compact summary"),
        (
            3050000,
            1480000,
            1800000,
            2750000,
            "Wide",
            "This description intentionally spans wider width",
        ),
        (6200000, 1515000, 1320000, 1700000, "Medium", "Balanced note"),
    ]

    for left, top, title_w, desc_w, title, desc in callouts:
        add_textbox(
            slide,
            left=left,
            top=top,
            width=title_w,
            height=340000,
            text=title,
            font_size=22,
        )
        add_textbox(
            slide,
            left=left,
            top=top + 390000,
            width=desc_w,
            height=380000,
            text=desc,
            font_size=18,
        )

    filename = "pptx_callout_blocks_mixed_widths.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_two_note_clusters(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Observations")

    add_textbox(
        slide,
        left=900000,
        top=1300000,
        width=1800000,
        height=340000,
        text="Latency",
        font_size=22,
    )
    add_textbox(
        slide,
        left=900000,
        top=1670000,
        width=3000000,
        height=370000,
        text="Improved after cache warm-up.",
        font_size=17,
    )

    add_textbox(
        slide,
        left=5600000,
        top=3150000,
        width=1850000,
        height=340000,
        text="Coverage",
        font_size=22,
    )
    add_textbox(
        slide,
        left=5600000,
        top=3520000,
        width=2600000,
        height=370000,
        text="Best on English queries.",
        font_size=17,
    )

    filename = "pptx_two_note_clusters.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_caption_scatter_one_real_pair(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Labels")

    labels = [
        (1050000, 1350000, "Search"),
        (4900000, 1450000, "Ranking"),
        (2400000, 2450000, "Safety"),
        (6000000, 2500000, "Vision"),
        (6150000, 2860000, "Fast setup"),
        (1350000, 3450000, "Internal only"),
    ]

    for left, top, text in labels:
        add_textbox(
            slide,
            left=left,
            top=top,
            width=1750000,
            height=340000,
            text=text,
            font_size=20,
        )

    filename = "pptx_caption_scatter_one_real_pair.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_caption_scatter_two_real_pairs(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Scatter Pairs")

    items = [
        (850000, 1300000, 1700000, "Alpha", 21),
        (850000, 1670000, 2550000, "Alpha detail", 17),
        (4500000, 1480000, 1700000, "Beta", 21),
        (4500000, 1850000, 2550000, "Beta detail", 17),
        (2600000, 2420000, 2000000, "Orphan tag", 20),
        (5850000, 2740000, 2400000, "Loose note", 18),
    ]

    for left, top, width, text, size in items:
        add_textbox(
            slide,
            left=left,
            top=top,
            width=width,
            height=340000,
            text=text,
            font_size=size,
        )

    filename = "pptx_caption_scatter_two_real_pairs.pptx"
    prs.save(output_dir / filename)
    return filename


def build_pptx_caption_scatter_pair_plus_footer_note(output_dir: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Scatter With Footer")

    top_region = [
        (1000000, 1350000, 1700000, "Signal", 21),
        (1000000, 1720000, 2300000, "Signal detail", 17),
        (4700000, 1490000, 1700000, "Noise", 21),
        (4700000, 1860000, 2300000, "Noise detail", 17),
        (2950000, 2520000, 2200000, "Standalone", 19),
    ]

    for left, top, width, text, size in top_region:
        add_textbox(
            slide,
            left=left,
            top=top,
            width=width,
            height=340000,
            text=text,
            font_size=size,
        )

    add_textbox(
        slide,
        left=900000,
        top=4720000,
        width=7300000,
        height=360000,
        text="Footer note: validate boundary handling.",
        font_size=16,
    )

    filename = "pptx_caption_scatter_pair_plus_footer_note.pptx"
    prs.save(output_dir / filename)
    return filename


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate PPTX caption-like / note-cluster regression samples.",
    )
    parser.add_argument(
        "output_dir",
        help="Directory to write generated .pptx files (example: samples/pptx).",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()
    if Presentation is None or Pt is None:
        print("Missing dependency: python-pptx. Install with: pip install python-pptx")
        sys.exit(1)

    out_dir = Path(args.output_dir)
    if not out_dir.is_absolute():
        out_dir = (ROOT / out_dir).resolve()

    ensure_dir(out_dir)

    generated = [
        build_pptx_callout_blocks_basic(out_dir),
        build_pptx_callout_blocks_row_jitter(out_dir),
        build_pptx_callout_blocks_mixed_widths(out_dir),
        build_pptx_two_note_clusters(out_dir),
        build_pptx_caption_scatter_one_real_pair(out_dir),
        build_pptx_caption_scatter_two_real_pairs(out_dir),
        build_pptx_caption_scatter_pair_plus_footer_note(out_dir),
    ]

    for name in generated:
        print(f"Generated: {name}")
