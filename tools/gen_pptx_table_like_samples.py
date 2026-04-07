#!/usr/bin/env python3
"""Generate table-like PPTX sample files and matching markdown expectations.

Usage:
    python tools/gen_pptx_table_like_samples.py <output_dir>

Dependency:
    pip install python-pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Callable, Iterable

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


MD_CONTENT: dict[str, str] = {
    "pptx_table_like_strong_2x3.md": """## Slide 1

### Result Matrix

Metric

Value

Latency

120ms

Accuracy

91%
""",
    "pptx_table_like_strong_3x3_header.md": """## Slide 1

### Quarter Overview

Quarter

Revenue

Profit

Q1

120

40

Q2

140

55
""",
    "pptx_table_like_local_edge_cell.md": """## Slide 1

### Local Metrics

CPU

65%

Memory

72%

Snapshot taken this morning.
""",
    "pptx_table_like_negative_cards_2x2.md": """## Slide 1

### Capabilities

Search

Fast retrieval

Vision

Image parsing

Speech

Audio input

Agents

Task execution
""",
    "pptx_table_like_negative_two_column_explainer.md": """## Slide 1

### Launch Stages

Plan

Set scope and owners

Build

Implement and test

Launch

Ship and monitor
""",
    "pptx_table_like_local_with_side_note.md": """## Slide 1

### Benchmark Snapshot

Model

Score

A

91

Best result highlighted.
""",
}


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold


def add_title(slide, text: str) -> None:
    add_textbox(
        slide,
        left=0.6,
        top=0.25,
        width=12.1,
        height=0.55,
        text=text,
        font_size=30,
        bold=True,
    )


def make_pptx_table_like_strong_2x3(out_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Result Matrix")

    left0, top0 = 2.1, 1.6
    col_w, row_h = 2.9, 0.8
    rows = [
        ("Metric", "Value"),
        ("Latency", "120ms"),
        ("Accuracy", "91%"),
    ]
    for r, (c1, c2) in enumerate(rows):
        y = top0 + r * row_h
        add_textbox(slide, left0, y, col_w, 0.55, c1, bold=(r == 0), font_size=20)
        add_textbox(slide, left0 + col_w + 0.25, y, col_w, 0.55, c2, bold=(r == 0), font_size=20)

    prs.save(out_path)


def make_pptx_table_like_strong_3x3_header(out_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Quarter Overview")

    left0, top0 = 1.3, 1.6
    col_w, row_h, gap = 2.8, 0.8, 0.2
    rows = [
        ("Quarter", "Revenue", "Profit"),
        ("Q1", "120", "40"),
        ("Q2", "140", "55"),
    ]
    for r, row in enumerate(rows):
        y = top0 + r * row_h
        for c, text in enumerate(row):
            x = left0 + c * (col_w + gap)
            add_textbox(slide, x, y, col_w, 0.55, text, bold=(r == 0), font_size=19)

    prs.save(out_path)


def make_pptx_table_like_local_edge_cell(out_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Local Metrics")

    left0, top0 = 2.4, 1.9
    col_w, row_h = 2.5, 0.9
    add_textbox(slide, left0, top0, col_w, 0.6, "CPU", bold=True, font_size=20)
    add_textbox(slide, left0 + col_w + 0.35, top0, col_w, 0.6, "65%", bold=True, font_size=20)
    add_textbox(slide, left0, top0 + row_h, col_w, 0.6, "Memory", font_size=20)
    add_textbox(
        slide,
        left0 + col_w + 0.42,
        top0 + row_h + 0.06,
        col_w - 0.25,
        0.58,
        "72%",
        font_size=20,
    )

    add_textbox(
        slide,
        left=2.4,
        top=4.15,
        width=6.8,
        height=0.55,
        text="Snapshot taken this morning.",
        font_size=16,
    )

    prs.save(out_path)


def make_pptx_table_like_negative_cards_2x2(out_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Capabilities")

    cards = [
        ("Search", "Fast retrieval"),
        ("Vision", "Image parsing"),
        ("Speech", "Audio input"),
        ("Agents", "Task execution"),
    ]
    start_left, start_top = 1.5, 1.8
    card_w, card_h = 4.6, 1.6
    col_gap, row_gap = 1.0, 0.9

    for i, (title, desc) in enumerate(cards):
        row, col = divmod(i, 2)
        x = start_left + col * (card_w + col_gap)
        y = start_top + row * (card_h + row_gap)
        add_textbox(slide, x, y, card_w, 0.48, title, font_size=21, bold=True)
        add_textbox(slide, x, y + 0.42, card_w, 0.44, desc, font_size=16)

    prs.save(out_path)


def make_pptx_table_like_negative_two_column_explainer(out_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Launch Stages")

    left_label, left_expl = 1.3, 3.6
    row_top, row_h = 1.9, 1.05
    pairs = [
        ("Plan", "Set scope and owners"),
        ("Build", "Implement and test"),
        ("Launch", "Ship and monitor"),
    ]

    for r, (label, explanation) in enumerate(pairs):
        y = row_top + r * row_h
        add_textbox(slide, left_label, y, 1.9, 0.58, label, font_size=20, bold=True)
        add_textbox(slide, left_expl, y, 6.9, 0.58, explanation, font_size=18)

    prs.save(out_path)


def make_pptx_table_like_local_with_side_note(out_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Benchmark Snapshot")

    left0, top0 = 2.3, 1.9
    col_w, row_h, gap = 2.6, 0.9, 0.32
    rows = [("Model", "Score"), ("A", "91")]
    for r, (c1, c2) in enumerate(rows):
        y = top0 + r * row_h
        add_textbox(slide, left0, y, col_w, 0.58, c1, font_size=20, bold=(r == 0))
        add_textbox(slide, left0 + col_w + gap, y, col_w, 0.58, c2, font_size=20, bold=(r == 0))

    add_textbox(
        slide,
        left=6.4,
        top=3.85,
        width=3.8,
        height=0.5,
        text="Best result highlighted.",
        font_size=15,
    )

    prs.save(out_path)


def generate_all(output_dir: Path) -> list[Path]:
    ensure_dir(output_dir)

    generators: Iterable[tuple[str, Callable[[Path], None]]] = [
        ("pptx_table_like_strong_2x3.pptx", make_pptx_table_like_strong_2x3),
        ("pptx_table_like_strong_3x3_header.pptx", make_pptx_table_like_strong_3x3_header),
        ("pptx_table_like_local_edge_cell.pptx", make_pptx_table_like_local_edge_cell),
        ("pptx_table_like_negative_cards_2x2.pptx", make_pptx_table_like_negative_cards_2x2),
        (
            "pptx_table_like_negative_two_column_explainer.pptx",
            make_pptx_table_like_negative_two_column_explainer,
        ),
        ("pptx_table_like_local_with_side_note.pptx", make_pptx_table_like_local_with_side_note),
    ]

    generated: list[Path] = []
    for filename, fn in generators:
        path = output_dir / filename
        fn(path)
        generated.append(path)

    return generated


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate table-like PPTX sample files.")
    parser.add_argument("output_dir", type=Path, help="Output directory for generated .pptx files")
    parser.add_argument(
        "--expected-md-dir",
        type=Path,
        default=Path("samples/expected/pptx"),
        help="Directory for generated expected .md files (default: samples/expected/pptx)",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    generated = generate_all(args.output_dir)
    ensure_dir(args.expected_md_dir)
    generated_md: list[Path] = []
    for md_name, content in MD_CONTENT.items():
        md_path = args.expected_md_dir / md_name
        md_path.write_text(content, encoding="utf-8")
        generated_md.append(md_path)

    print("Generated files:")
    for p in generated:
        print(f"- {p}")
    print("Generated expected markdown files:")
    for p in generated_md:
        print(f"- {p}")


if __name__ == "__main__":
    main()
